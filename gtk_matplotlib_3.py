import gi

gi.require_version('Gtk', '3.0')
from gi.repository import Gtk, GdkPixbuf

from matplotlib.backends.backend_gtk3agg import (
    FigureCanvasGTK3Agg as FigureCanvas
)
from matplotlib.figure import Figure
import pandas as pd

from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor


class MyWindow(Gtk.Window):

    def __init__(self):
        Gtk.Window.__init__(self, title="SPC Chart")
        self.set_default_size(800, 600)

        # SPC chart
        figure = self.gen_example_chart()

        box = Gtk.Box(orientation=Gtk.Orientation.VERTICAL)
        self.add(box)

        # frame
        b_bar = self.gen_button_bar(figure)
        box.pack_start(b_bar, False, True, 0)

        # plot area
        p_area = self.gen_plot_area(figure)
        box.pack_start(p_area, True, True, 0)

    def gen_button_bar(self, figure):
        frame = Gtk.Frame()

        hbox = Gtk.Box(orientation=Gtk.Orientation.HORIZONTAL)
        frame.add(hbox)

        pixbuf = GdkPixbuf.Pixbuf.new_from_file('powerpoint-128.png')
        pixbuf = pixbuf.scale_simple(32, 32, GdkPixbuf.InterpType.BILINEAR)

        but = Gtk.Button()
        but.add(Gtk.Image.new_from_pixbuf(pixbuf))
        but.connect("clicked", self.on_button_clicked, figure)
        hbox.pack_end(but, False, True, 0)

        return frame

    def gen_plot_area(self, figure):
        sw = Gtk.ScrolledWindow()
        sw.set_border_width(0)
        canvas = FigureCanvas(figure)  # a Gtk.DrawingArea
        canvas.set_size_request(800, 600)
        sw.add(canvas)
        return sw

    def gen_example_chart(self):
        # example dataframe
        df = pd.DataFrame({
            'Sample': list(range(1, 11)),
            'Y': [9.030, 8.810, 9.402, 8.664, 8.773, 8.774, 8.416, 9.101, 8.687, 8.767]
        })
        # SPC metrics
        spec_usl = 9.97
        spec_target = 8.70
        spec_lsl = 7.43

        # spc chart
        fig = Figure(dpi=100)
        splot = fig.add_subplot(111, title="SPC Chart Example", ylabel='Value')
        splot.grid(True)

        # horizontal lines
        splot.axhline(y=spec_usl, linewidth=1, color='red', label='USL')
        splot.axhline(y=spec_target, linewidth=1, color='blue', label='Target')
        splot.axhline(y=spec_lsl, linewidth=1, color='red', label='LSL')

        # trend
        splot.plot(df['Sample'], df['Y'], color="gray", marker="o", markersize=10)

        # text label
        x_label = splot.get_xlim()[1]
        splot.text(x_label, spec_usl, " USL", color="red")
        splot.text(x_label, spec_target, " Target", color="blue")
        splot.text(x_label, spec_lsl, " LSL", color="red")

        return fig

    def on_button_clicked(self, button, figure):
        image_path = "./chart.png"
        template_path = "./template.pptx"
        save_path = "./output.pptx"

        # create PNG file of plots
        figure.savefig(image_path)

        # create PowerPoint files
        self.gen_ppt(image_path, template_path, save_path)

        # complete messages
        dialog = Gtk.MessageDialog(parent=self,
                                   flags=0,
                                   message_type=Gtk.MessageType.INFO,
                                   buttons=Gtk.ButtonsType.OK,
                                   text="generated PowerPoint file in " + save_path + ".")
        dialog.run()
        dialog.destroy()

    def gen_ppt(self, image_path, template_path, save_path):
        # insert empty slide
        presentation = Presentation(template_path)
        # refer layout from original master
        title_slide_layout = presentation.slide_layouts[3]
        slide = presentation.slides.add_slide(title_slide_layout)
        shapes = slide.shapes
        # slide title
        slide_title = "SPC chart (Example)"
        shapes.title.text = slide_title

        # ---------------------------------------------------
        # insert image
        # ---------------------------------------------------
        # insert position
        pic_left = Cm(0.5)
        pic_top = Cm(3.5)
        # image height
        pic_height = Cm(15)
        slide.shapes.add_picture(image_path, pic_left, pic_top, height=pic_height)

        # ---------------------------------------------------
        # insert text box
        # ---------------------------------------------------
        text_left = Cm(20.73)
        text_top = Cm(3.63)
        text_height = Cm(14.88)
        text_width = Cm(11.63)
        text_box = slide.shapes.add_textbox(text_left, text_top, text_width, text_height)

        sample_str = "Test Message"
        text_box.text = sample_str

        text_font = 20
        text_box.text_frame.add_paragraph().font.size = Pt(text_font)

        text_box.fill.solid()
        color = RGBColor(255, 255, 240)
        text_box.fill.fore_color.rgb = color

        # ---------------------------------------------------------------------
        # save PowerPoint file
        # ---------------------------------------------------------------------
        presentation.save(save_path)


win = MyWindow()
win.connect("destroy", Gtk.main_quit)
win.show_all()
Gtk.main()
